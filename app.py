from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path
import json

import spacy

INPUT_DIR = Path("input")
OUTPUT_DIR = Path("output")
JSON_FILE = OUTPUT_DIR / "graph.json"
MODEL = "en_core_web_sm"

MIN_MENTIONS = 1

MERGE_ALIASES = True

NAMED_LABELS = {
    "PERSON": "People, including fictional",
    "NORP": "Nationalities, religious or political groups",
    "FAC": "Buildings, airports, highways, bridges, etc.",
    "ORG": "Companies, agencies, institutions",
    "GPE": "Countries, cities, states",
    "LOC": "Non-GPE locations: mountains, bodies of water",
    "PRODUCT": "Objects, vehicles, foods, etc. (not services)",
    "EVENT": "Named hurricanes, battles, wars, sports events",
    "WORK_OF_ART": "Titles of books, songs, etc.",
    "LAW": "Named documents made into laws",
    "LANGUAGE": "Any named language",
}


def read_pdf(path):
    import fitz

    parts = []
    with fitz.open(path) as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n".join(parts)


def _iter_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def read_pptx(path):
    from pptx import Presentation

    parts = []
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(notes)
    return "\n".join(parts)


def read_txt(path):
    try:
        return path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


READERS = {
    ".pdf": (read_pdf, "pdf"),
    ".pptx": (read_pptx, "pptx"),
    ".txt": (read_txt, "txt"),
}


def normalize(text):
    return " ".join(text.split())


def chunks(text, size):
    if len(text) <= size:
        yield text
        return
    start, n = 0, len(text)
    while start < n:
        end = min(start + size, n)
        if end < n:
            split = text.rfind("\n", start, end)
            if split > start:
                end = split
        yield text[start:end]
        start = end


def entities_in(text, nlp):
    for doc in nlp.pipe(chunks(text, nlp.max_length - 1)):
        for ent in doc.ents:
            if ent.label_ not in NAMED_LABELS:
                continue
            name = normalize(ent.text)
            if not name:
                continue
            yield name, (name.lower(), ent.label_), ent.label_


def _contains_tokens(short, long):
    n = len(short)
    return any(long[i:i + n] == short for i in range(len(long) - n + 1))


def merge_aliases(kept, ent_total):
    by_label = defaultdict(list)
    for key in kept:
        by_label[key[1]].append(key)

    canonical = {key: key for key in kept}
    for keys in by_label.values():
        token_index = defaultdict(list)
        for key in keys:
            for tok in set(key[0].split()):
                token_index[tok].append(key)
        for short in keys:
            short_tokens = short[0].split()
            best, best_total = None, -1
            for long in token_index.get(short_tokens[0], ()):
                long_tokens = long[0].split()
                if len(long_tokens) <= len(short_tokens):
                    continue
                if _contains_tokens(short_tokens, long_tokens) and (
                    ent_total[long] > best_total
                ):
                    best, best_total = long, ent_total[long]
            if best is not None:
                canonical[short] = best

    def resolve(key):
        while canonical[key] != key:
            key = canonical[key]
        return key

    return {key: resolve(key) for key in kept}


def main():
    files = sorted(
        p for p in INPUT_DIR.glob("*") if p.suffix.lower() in READERS
    )
    skipped = sorted(
        p.name for p in INPUT_DIR.glob("*")
        if p.is_file() and p.suffix.lower() not in READERS
    )
    if skipped:
        print(f"Skipping unsupported file(s): {', '.join(skipped)}")
        if any(name.lower().endswith(".ppt") for name in skipped):
            print(
                "  Legacy .ppt (binary PowerPoint) can't be read; "
                "convert it to .pptx first."
            )
    if not files:
        raise SystemExit(
            f"No .pdf, .pptx or .txt files found in {INPUT_DIR}/"
        )

    try:
        nlp = spacy.load(
            MODEL, disable=["tagger", "parser", "attribute_ruler", "lemmatizer"]
        )
    except OSError:
        raise SystemExit(
            f"spaCy model '{MODEL}' is not installed. Install it with:\n"
            f"    pip install -r requirements.txt\n"
            f"or: python -m spacy download {MODEL}"
        )

    edge = Counter()
    ent_total = Counter()
    ent_docs = defaultdict(set)
    ent_surface = defaultdict(Counter)
    doc_formats = {}

    total = len(files)
    for i, path in enumerate(files, 1):
        reader, fmt = READERS[path.suffix.lower()]
        doc_formats[path.name] = fmt
        print(f"[{i}/{total}] {path.name}", flush=True)
        try:
            text = reader(path)
        except Exception as exc:
            print(f"          ! could not read: {exc}", flush=True)
            continue

        seen = set()
        for display, key, _label in entities_in(text, nlp):
            edge[(path.name, key)] += 1
            ent_total[key] += 1
            ent_docs[key].add(path.name)
            ent_surface[key][display] += 1
            seen.add(key)

        print(f"          {len(seen)} distinct entities", flush=True)

    kept = {k for k, count in ent_total.items() if count >= MIN_MENTIONS}
    if not kept:
        raise SystemExit("No entities found in the input documents.")

    if MERGE_ALIASES:
        alias = merge_aliases(kept, ent_total)
        if any(alias[k] != k for k in alias):
            merged_edge = Counter()
            for (doc_name, key), c in edge.items():
                if key in kept:
                    merged_edge[(doc_name, alias[key])] += c
            merged_total = Counter()
            merged_docs = defaultdict(set)
            merged_surface = defaultdict(Counter)
            for key in kept:
                canon = alias[key]
                merged_total[canon] += ent_total[key]
                merged_docs[canon] |= ent_docs[key]
                merged_surface[canon].update(ent_surface[key])
            edge, ent_total = merged_edge, merged_total
            ent_docs, ent_surface = merged_docs, merged_surface
            kept = set(alias.values())
            print(f"Merged aliases: down to {len(kept)} distinct entities.")

    nodes = []

    doc_mentions = Counter()
    doc_distinct = Counter()
    for (doc_name, key), count in edge.items():
        if key in kept:
            doc_mentions[doc_name] += count
            doc_distinct[doc_name] += 1
    for name in sorted(doc_mentions):
        nodes.append({
            "id": f"doc::{name}",
            "type": "document",
            "label": name,
            "format": doc_formats[name],
            "entityCount": doc_distinct[name],
            "mentionCount": doc_mentions[name],
        })

    def ent_id(key):
        name_lc, label = key
        return f"ent::{label}::{name_lc}"

    for key in sorted(kept, key=lambda k: (-ent_total[k], k)):
        _name_lc, label = key
        display = ent_surface[key].most_common(1)[0][0]
        nodes.append({
            "id": ent_id(key),
            "type": "entity",
            "label": display,
            "entityType": label,
            "entityTypeLabel": NAMED_LABELS[label],
            "docCount": len(ent_docs[key]),
            "mentionCount": ent_total[key],
        })

    links = []
    for (doc_name, key), count in sorted(edge.items()):
        if key not in kept:
            continue
        links.append({
            "source": f"doc::{doc_name}",
            "target": ent_id(key),
            "value": count,
        })

    graph = {
        "meta": {
            "generated": datetime.now(timezone.utc).isoformat(timespec="seconds"),
            "documents": len(doc_mentions),
            "entities": len(kept),
            "links": len(links),
            "entityTypes": NAMED_LABELS,
        },
        "nodes": nodes,
        "links": links,
    }

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    JSON_FILE.write_text(json.dumps(graph, indent=2), encoding="utf-8")

    print(
        f"\nWrote {len(nodes)} nodes "
        f"({len(doc_mentions)} documents, {len(kept)} entities) "
        f"and {len(links)} links."
    )
    print(f"  data: {JSON_FILE}")
    print(
        "  View it by serving the repo root over http (e.g. VS Code Live "
        f"Server) and opening index.html; it fetches {JSON_FILE} at runtime."
    )


if __name__ == "__main__":
    main()
